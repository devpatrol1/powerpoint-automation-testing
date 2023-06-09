import collections 
import collections.abc
from pptx import Presentation
from PIL import Image
import os



def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]
 
    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size
 
    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width
 
    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)
 
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
 
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side



prs = Presentation()

layout8 = prs.slide_layouts[8]
slide = prs.slides.add_slide(layout8)
 
title = slide.shapes.title.text = "This is Powerpoint"
sub = slide.placeholders[2].text = "Python has the power"
_add_image(slide,1,"images/butterfly5.png")
 
prs.save("ppt_files/MyPresentation.pptx")
#os.startfile("ppt_files/MyPresentation.pptx")
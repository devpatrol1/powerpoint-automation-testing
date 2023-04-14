import collections 
import collections.abc
from pptx import Presentation
import os
 
prs = Presentation()
 
class MySlide:
    def __init__(self, data):
        self.layout = prs.slide_layouts[data[3]]
        self.slide=prs.slides.add_slide(self.layout)
        self.title=self.slide.shapes.title
        self.title.text=data[0]
        self.subtitle=self.slide.placeholders[1]
        self.subtitle.text=data[1]
 
        for shape in self.slide.placeholders:
            shape.placeholder_format.idx,
            shape.placeholder_format.type,
            shape.name
        if data[2] != "":
        	self.img = self.slide.placeholders[1].insert_picture(data[2])
 
slides = [
    ["USA Weather",       #data[0]
     "Subtitle(Bullet)",
     "images/butterfly5.png",
     8],
    ["Malaysia Weather",       #data[0]
     "Content(Bullet)",
     "",
     3],
    ["China Weather",       #data[0]
     "This is a brown Fox",
     "",
     3]
]
 
for each_slide in slides:
    MySlide(each_slide)
 
prs.save("ppt_files/stack.pptx")
#os.startfile("ppt_files/stack.pptx")
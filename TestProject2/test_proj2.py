import collections 
import collections.abc
from pptx import Presentation
from PIL import Image, ImageOps
import os, sys


def get_placeholder_info(prs):
    slide_layouts = prs.slide_master.slide_layouts
    for num in range(len(slide_layouts)):
        for placeholder in prs.slide_layouts[num].placeholders:
            print('%d %s' % (placeholder.placeholder_format.idx, placeholder.name))


def pad_img(path):
    img_dir = os.listdir(path)
    for item in img_dir:
        if os.path.isfile(path+item):
            img = Image.open(path+item)
            right = 200
            left = 200
            top = 100
            bottom = 100       
            width, height = img.size
            new_width = width + right + left
            new_height = height + top + bottom            
            new_img = Image.new(img.mode, (new_width, new_height), (255, 255, 255))
            #place old image on new image           
            new_img.paste(img, (left, top))
            new_img.save(path + item)


def add_slide(prs, layout, title):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    return slide


def update_topic_slide(tslide, num):
    subtitle_placeholder = tslide.placeholders[1]
    subtitle_placeholder.text = 'Structured Streaming Method #'+ num


def update_data_slide(root, topic, dslide):
    text1_placeholder = dslide.placeholders[1]
    text2_placeholder = dslide.placeholders[3]
    chart1_placeholder = dslide.placeholders[13]
    chart2_placeholder = dslide.placeholders[15]
    data1_placeholder = dslide.placeholders[14]
    data2_placeholder = dslide.placeholders[16]
    text1_placeholder.text = 'Example of Sub-Header #1'
    text2_placeholder.text = 'Example of Sub-Header #2'
    chart1img = chart1_placeholder.insert_picture(root + topic + '/chart1.png')
    chart2img = chart2_placeholder.insert_picture(root + topic + '/chart2.png')
    data1img = data1_placeholder.insert_picture(root + topic + '/data1.png')
    data2img = data2_placeholder.insert_picture(root + topic + '/data2.png')
    return chart1img, chart2img, data1img, data2img


def update_process_slide(root, topic, pslide):
    pic1_placeholder = pslide.placeholders[13]
    pic2_placeholder = pslide.placeholders[15]
    pic1img = pic1_placeholder.insert_picture(root + topic + '/pic1.png')
    pic2img = pic2_placeholder.insert_picture(root + topic + '/pic2.png')
    return pic1img, pic2img




prs = Presentation('test_ppt2_template.pptx')
#get_placeholder_info(prs) #Only needed to find out placeholder names & indexes

root ='./presentation_info/'
folder_list = [ item for item in os.listdir(root) if os.path.isdir(os.path.join(root, item)) ]
all_slides = prs.slides

#Pad images for insertion in slides 
for x in range (len(folder_list)):
    topic = folder_list[x]
    path = root + topic + '/'
    pad_img(path)
    
# Update existing intro slide
main_topic = 'Methods for Writing Continuous Applications'
intro_slide = all_slides[0]
intro_slide_title = intro_slide.shapes.title
intro_slide_title.text = 'Structured Streaming'
intro_slide_subtitle = intro_slide.placeholders[1]
intro_slide_subtitle.text = main_topic

# Add slides for each name in the folder list 
for x in range (len(folder_list)):
    topic = folder_list[x]
    # Create slide for topic title
    new_slide1_layout = prs.slide_layouts[1]
    new_slide1 = add_slide(prs, new_slide1_layout, topic)
    # Create slide for charts & data
    new_slide2_layout = prs.slide_layouts[2]
    new_slide2 = add_slide(prs, new_slide2_layout, topic)
    # Create slide for process
    new_slide3_layout = prs.slide_layouts[3]
    new_slide3 = add_slide(prs, new_slide3_layout, topic)   

#Update slides with subtitles and images by topic
num = 0
for x in range(1, len(all_slides), 3):
    update_topic_slide(all_slides[x], str(num+1))
    update_data_slide(root, folder_list[num], all_slides[x+1])
    update_process_slide(root, folder_list[num], all_slides[x+2])
    num += 1

prs.save('test_ppt2_finalcopy.pptx')


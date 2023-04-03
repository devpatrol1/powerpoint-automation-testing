import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches

img_path = 'butterfly5.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(0.5)
# pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(1)
height = Inches(1.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs.save('ppt_files/test4.pptx')